"""
presentation_importer.py — higher-level PDF/PPTX → Markdown importers.

Key differences from src.converters.*:
  • Images saved to <input_dir>/images/ with slide-aware naming
    (slide_N_img_M.png).
  • PPTX presenter notes emitted as <!-- notes: TEXT --> HTML comments.
  • ImportResult captures conversion statistics for CLI reporting.

Classes:
    ImportResult          — value object returned by every importer
    BaseImporter          — ABC defining the import_file interface
    PPTXImporter          — PPTX → Markdown (requires python-pptx)
    PDFImporter           — PDF  → Markdown (requires pdfplumber, Pillow)
    PresentationImporter  — factory: selects the right importer by extension

SRP: each class has one reason to change.
OCP: new format support is added by subclassing BaseImporter.
DIP: import_presentation.py depends on BaseImporter, not concrete classes.
"""

from abc import ABC, abstractmethod
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple

from src.converters.layout_analyzer import LayoutAnalyzer, ShapeRegion
from src.grid_builder import GridBuilder


# ---------------------------------------------------------------------------
# Value object
# ---------------------------------------------------------------------------

@dataclass
class ImportResult:
    """Carries both the generated Markdown and conversion statistics."""
    md_content: str
    slide_count: int
    image_count: int
    output_path: Path


# ---------------------------------------------------------------------------
# Abstract base
# ---------------------------------------------------------------------------

class BaseImporter(ABC):
    """
    Abstract base for all presentation importers.

    OCP: new formats are added by subclassing — no existing code is touched.
    ISP: a single minimal method, a single concern.
    DIP: PresentationImporter factory and the CLI depend on this abstraction.
    """

    SLIDE_SEP = "\n\n---\n\n"

    @abstractmethod
    def import_file(
        self,
        input_path: Path,
        output_path: Optional[Path] = None,
    ) -> ImportResult:
        """
        Import a presentation file and write the resulting Markdown.

        Args:
            input_path:  Path to the .pptx or .pdf source file.
            output_path: Explicit output .md path, or None to write the file
                         next to the source with the same stem.

        Returns:
            ImportResult with Markdown content and conversion statistics.
        """

    # ------------------------------------------------------------------
    # Shared helpers
    # ------------------------------------------------------------------

    def _resolve_paths(
        self,
        input_path: Path,
        output_path: Optional[Path],
    ) -> Tuple[Path, Path]:
        """Return (output_md_path, images_dir)."""
        out = output_path or input_path.with_suffix(".md")
        images_dir = input_path.parent / "images"
        return out, images_dir


# ---------------------------------------------------------------------------
# PPTX importer
# ---------------------------------------------------------------------------

class PPTXImporter(BaseImporter):
    """
    Converts a .pptx file to md-reveal-wrapper Markdown.

    Requires: python-pptx

    Differences from PPTXConverter:
      • Images are saved to images/slide_N_img_M.png next to the source file.
      • Presenter notes are appended as <!-- notes: TEXT --> per slide.
      • Returns an ImportResult with slide/image statistics.

    Text, table, and heading logic is delegated to PPTXConverter to avoid
    duplication, while image handling is fully owned here.
    """

    # Placeholder indices treated as "title" by PowerPoint convention.
    _TITLE_IDX = {0, 1}

    def import_file(
        self,
        input_path: Path,
        output_path: Optional[Path] = None,
    ) -> ImportResult:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required: pip install python-pptx"
            ) from exc

        from src.converters.pptx_converter import PPTXConverter
        _conv = PPTXConverter()

        out_path, images_dir = self._resolve_paths(input_path, output_path)
        prs = Presentation(str(input_path))
        analyzer = LayoutAnalyzer()

        slide_mds: List[str] = []
        total_images = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            img_counter = [0]
            slide_md = self._convert_slide(
                slide, prs, images_dir, slide_num, img_counter, analyzer, _conv
            )
            notes = self._extract_notes(slide)
            if notes:
                slide_md = slide_md.rstrip() + f"\n\n<!-- notes: {notes} -->"

            slide_mds.append(slide_md)
            total_images += img_counter[0]

        md_content = self.SLIDE_SEP.join(slide_mds)
        out_path.write_text(md_content, encoding="utf-8")

        return ImportResult(
            md_content=md_content,
            slide_count=len(slide_mds),
            image_count=total_images,
            output_path=out_path,
        )

    # ------------------------------------------------------------------
    # Per-slide conversion
    # ------------------------------------------------------------------

    def _convert_slide(
        self,
        slide,
        prs,
        images_dir: Path,
        slide_num: int,
        img_counter: List[int],
        analyzer: LayoutAnalyzer,
        conv,
    ) -> str:
        from pptx.util import Emu

        sw = float(prs.slide_width or Emu(9144000))
        sh = float(prs.slide_height or Emu(6858000))

        # Sort shapes in natural reading order (top-to-bottom, left-to-right).
        shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))

        title_md = ""
        regions: List[ShapeRegion] = []

        for shape in shapes:
            is_title = self._is_title(shape)
            content_md = self._shape_to_markdown(
                shape, images_dir, slide_num, img_counter, conv
            )
            if not content_md.strip():
                continue

            x = (shape.left or 0) / sw
            y = (shape.top or 0) / sh
            w = (shape.width or sw) / sw
            h = (shape.height or sh) / sh

            if is_title:
                title_md = content_md.strip().lstrip("#").strip()
            else:
                regions.append(ShapeRegion(content_md=content_md, x=x, y=y, w=w, h=h))

        layout = analyzer.analyze(regions)
        if title_md:
            layout.title = title_md

        return analyzer.to_markdown(layout)

    # ------------------------------------------------------------------
    # Shape dispatch
    # ------------------------------------------------------------------

    def _shape_to_markdown(
        self,
        shape,
        images_dir: Path,
        slide_num: int,
        img_counter: List[int],
        conv,
    ) -> str:
        if hasattr(shape, "table"):
            return conv._table_to_markdown(shape.table)

        if hasattr(shape, "text_frame"):
            return conv._text_frame_to_markdown(shape.text_frame)

        if conv._is_image(shape):
            return self._save_image(shape, images_dir, slide_num, img_counter)

        return ""

    def _is_title(self, shape) -> bool:
        try:
            return (
                shape.placeholder_format is not None
                and shape.placeholder_format.idx in self._TITLE_IDX
            )
        except Exception:
            return False

    def _save_image(
        self,
        shape,
        images_dir: Path,
        slide_num: int,
        img_counter: List[int],
    ) -> str:
        try:
            img_counter[0] += 1
            images_dir.mkdir(parents=True, exist_ok=True)
            img_name = f"slide_{slide_num}_img_{img_counter[0]}.png"
            (images_dir / img_name).write_bytes(shape.image.blob)
            return f"![](images/{img_name})"
        except Exception:
            return ""

    def _extract_notes(self, slide) -> str:
        try:
            text = slide.notes_slide.notes_text_frame.text
            return text.strip()
        except Exception:
            return ""


# ---------------------------------------------------------------------------
# PDF importer
# ---------------------------------------------------------------------------

class PDFImporter(BaseImporter):
    """
    Converts a .pdf file to md-reveal-wrapper Markdown.

    Requires: pdfplumber, Pillow

    Differences from PDFConverter:
      • Images are saved to images/slide_N_img_M.png next to the source file.
      • Returns an ImportResult with slide/image statistics.

    Word-grouping and line-classification logic is delegated to PDFConverter
    to avoid duplication; image handling is fully owned here.
    """

    def import_file(
        self,
        input_path: Path,
        output_path: Optional[Path] = None,
    ) -> ImportResult:
        try:
            import pdfplumber
        except ImportError as exc:
            raise ImportError(
                "pdfplumber is required: pip install pdfplumber"
            ) from exc

        from src.converters.pdf_converter import PDFConverter
        _conv = PDFConverter()

        out_path, images_dir = self._resolve_paths(input_path, output_path)
        analyzer = LayoutAnalyzer()

        slide_mds: List[str] = []
        total_images = 0

        with pdfplumber.open(str(input_path)) as pdf:
            for slide_num, page in enumerate(pdf.pages, start=1):
                img_counter = [0]
                slide_md = self._convert_page(
                    page, images_dir, slide_num, img_counter, analyzer, _conv
                )
                slide_mds.append(slide_md)
                total_images += img_counter[0]

        md_content = self.SLIDE_SEP.join(slide_mds)
        out_path.write_text(md_content, encoding="utf-8")

        return ImportResult(
            md_content=md_content,
            slide_count=len(slide_mds),
            image_count=total_images,
            output_path=out_path,
        )

    # ------------------------------------------------------------------
    # Per-page conversion
    # ------------------------------------------------------------------

    def _convert_page(
        self,
        page,
        images_dir: Path,
        slide_num: int,
        img_counter: List[int],
        analyzer: LayoutAnalyzer,
        conv,
    ) -> str:
        import statistics

        pw = page.width or 1.0
        ph = page.height or 1.0

        words = page.extract_words(extra_attrs=["size", "fontname"])
        lines = conv._group_words_into_lines(words)
        blocks = conv._group_lines_into_blocks(lines)

        all_sizes = [w.get("size", 12) for w in words if w.get("size")]
        median_size = statistics.median(all_sizes) if all_sizes else 12.0

        regions: List[ShapeRegion] = []

        for block in blocks:
            content_md, bx, by, bw, bh = conv._block_to_markdown(
                block, median_size, pw, ph
            )
            if content_md.strip():
                regions.append(ShapeRegion(
                    content_md=content_md, x=bx, y=by, w=bw, h=bh
                ))

        for img_info in (page.images or []):
            img_md = self._extract_image(
                page, img_info, images_dir, slide_num, img_counter, pw, ph
            )
            if img_md:
                x0 = img_info.get("x0", 0)
                top = img_info.get("top", 0)
                x1 = img_info.get("x1", pw)
                bottom = img_info.get("bottom", ph)
                regions.append(ShapeRegion(
                    content_md=img_md,
                    x=x0 / pw,
                    y=top / ph,
                    w=(x1 - x0) / pw,
                    h=(bottom - top) / ph,
                ))

        layout = analyzer.analyze(regions)
        return analyzer.to_markdown(layout)

    # ------------------------------------------------------------------
    # Image extraction
    # ------------------------------------------------------------------

    def _extract_image(
        self,
        page,
        img_info: dict,
        images_dir: Path,
        slide_num: int,
        img_counter: List[int],
        pw: float,
        ph: float,
    ) -> str:
        try:
            img_counter[0] += 1
            images_dir.mkdir(parents=True, exist_ok=True)

            x0 = img_info.get("x0", 0)
            top = img_info.get("top", 0)
            x1 = img_info.get("x1", pw)
            bottom = img_info.get("bottom", ph)

            cropped = page.crop((x0, top, x1, bottom))
            pil_img = cropped.to_image(resolution=150)

            img_name = f"slide_{slide_num}_img_{img_counter[0]}.png"
            pil_img.save(str(images_dir / img_name))
            return f"![](images/{img_name})"
        except Exception:
            return ""


# ---------------------------------------------------------------------------
# Factory
# ---------------------------------------------------------------------------

class PresentationImporter:
    """
    Factory that selects the correct BaseImporter by file extension.

    DIP: callers depend on this factory and on BaseImporter, never on
         concrete importer classes directly.
    """

    _REGISTRY = {
        ".pptx": PPTXImporter,
        ".ppt":  PPTXImporter,
        ".pdf":  PDFImporter,
    }

    @classmethod
    def for_file(cls, path: Path) -> BaseImporter:
        """
        Return the appropriate importer for *path*.

        Raises:
            ValueError: If the file extension is not supported.
        """
        ext = path.suffix.lower()
        importer_cls = cls._REGISTRY.get(ext)
        if importer_cls is None:
            supported = ", ".join(cls._REGISTRY)
            raise ValueError(
                f"Unsupported format '{ext}'. Supported: {supported}"
            )
        return importer_cls()
