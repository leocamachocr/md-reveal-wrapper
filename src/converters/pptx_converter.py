import re
from pathlib import Path
from typing import List, Optional, Tuple

from src.converters.base import PresentationConverter
from src.converters.layout_analyzer import LayoutAnalyzer, ShapeRegion


class PPTXConverter(PresentationConverter):
    """
    Converts a .pptx file to md-reveal-wrapper Markdown.

    Requires: python-pptx

    SRP: sole responsibility is PPTX → Markdown translation.
    OCP: shape-type handlers can be extended without changing the pipeline.
    """

    # Placeholder indices considered "title" by PowerPoint convention
    TITLE_PLACEHOLDER_IDX = {0, 1}

    def convert(self, input_path: Path) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Emu
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required: pip install python-pptx"
            ) from exc

        prs = Presentation(str(input_path))
        assets_dir = input_path.parent / f"{input_path.stem}_assets"
        analyzer = LayoutAnalyzer()

        slide_mds: List[str] = []
        img_counter = [0]

        for slide in prs.slides:
            slide_md = self._convert_slide(
                slide, prs, assets_dir, img_counter, analyzer
            )
            slide_mds.append(slide_md)

        return "\n\n---\n\n".join(slide_mds)

    # ------------------------------------------------------------------
    # Per-slide conversion
    # ------------------------------------------------------------------

    def _convert_slide(self, slide, prs, assets_dir: Path, img_counter: List[int], analyzer: LayoutAnalyzer) -> str:
        from pptx.util import Emu

        slide_width = prs.slide_width or Emu(9144000)
        slide_height = prs.slide_height or Emu(6858000)
        sw = float(slide_width)
        sh = float(slide_height)

        shapes = list(slide.shapes)
        # Sort by (top, left) so reading order is natural
        shapes.sort(key=lambda s: (s.top or 0, s.left or 0))

        title_md = ""
        regions: List[ShapeRegion] = []

        for shape in shapes:
            is_title = self._is_title_placeholder(shape)
            content_md = self._shape_to_markdown(shape, assets_dir, img_counter)
            if not content_md.strip():
                continue

            x = (shape.left or 0) / sw
            y = (shape.top or 0) / sh
            w = (shape.width or sw) / sw
            h = (shape.height or sh) / sh

            if is_title:
                # Strip leading # added by heading detection; we'll re-add
                title_md = content_md.strip().lstrip("#").strip()
            else:
                regions.append(ShapeRegion(content_md=content_md, x=x, y=y, w=w, h=h))

        layout = analyzer.analyze(regions)
        # Override title from explicit title placeholder if found
        if title_md:
            layout.title = title_md

        return analyzer.to_markdown(layout)

    # ------------------------------------------------------------------
    # Shape → Markdown helpers
    # ------------------------------------------------------------------

    def _is_title_placeholder(self, shape) -> bool:
        try:
            return shape.placeholder_format is not None and \
                   shape.placeholder_format.idx in self.TITLE_PLACEHOLDER_IDX
        except Exception:
            return False

    def _shape_to_markdown(self, shape, assets_dir: Path, img_counter: List[int]) -> str:
        """Dispatch to the appropriate handler based on shape type."""
        # Callout detection by shape name
        callout_type = self._detect_callout(shape)

        if hasattr(shape, "table"):
            return self._table_to_markdown(shape.table)

        if hasattr(shape, "text_frame"):
            md = self._text_frame_to_markdown(shape.text_frame)
            if callout_type and md.strip():
                first_line, *rest = md.strip().split("\n", 1)
                body = rest[0] if rest else ""
                if body.strip():
                    return f"> [{callout_type}] {first_line.strip()}\n> {body.strip()}"
                return f"> [{callout_type}] {first_line.strip()}"
            return md

        if self._is_image(shape):
            return self._image_to_markdown(shape, assets_dir, img_counter)

        return ""

    def _detect_callout(self, shape) -> Optional[str]:
        """Check shape name or first text token for info/warning/tip."""
        name_lower = (shape.name or "").lower()
        for tag in ("info", "warning", "tip"):
            if tag in name_lower:
                return tag

        try:
            first_text = shape.text_frame.text.strip().lower()
            for tag in ("info", "warning", "tip"):
                if first_text.startswith(tag):
                    return tag
        except Exception:
            pass
        return None

    def _text_frame_to_markdown(self, text_frame) -> str:
        lines: List[str] = []
        for para in text_frame.paragraphs:
            md_line = self._paragraph_to_markdown(para)
            if md_line is not None:
                lines.append(md_line)
        return "\n".join(lines)

    def _paragraph_to_markdown(self, para) -> Optional[str]:
        runs = para.runs
        if not runs:
            raw_text = para.text
            if not raw_text.strip():
                return ""
            return raw_text

        # Detect list item via paragraph level
        level = para.level or 0
        indent = "  " * level

        # Detect heading via bold + large font on the first run
        if self._is_heading_para(para):
            text = "".join(r.text for r in runs)
            depth = self._heading_depth(para)
            return f"{'#' * depth} {text.strip()}"

        # Build inline-formatted text
        text = self._runs_to_inline(runs)
        if not text.strip():
            return ""

        # List item?
        if level > 0 or self._looks_like_list_item(para):
            return f"{indent}- {text.strip()}"

        return text.strip()

    def _is_heading_para(self, para) -> bool:
        runs = para.runs
        if not runs:
            return False
        first_run = runs[0]
        try:
            bold = first_run.font.bold
            size = first_run.font.size  # in EMU (914400 = 72 pt)
            if bold and size and size >= 914400 * 18:  # >= 18 pt
                return True
        except Exception:
            pass
        # Also check paragraph style name
        try:
            style_name = para._p.pPr.buNone  # crude check
        except Exception:
            pass
        return False

    def _heading_depth(self, para) -> int:
        try:
            size = para.runs[0].font.size or 0
            pt = size / 914400  # EMU → pt (actually 12700 EMU = 1pt for pptx)
            # pptx uses 12700 EMU per point
            pt = size / 12700
            if pt >= 36:
                return 1
            if pt >= 28:
                return 2
            return 3
        except Exception:
            return 2

    def _looks_like_list_item(self, para) -> bool:
        text = para.text.strip()
        return bool(re.match(r"^[\u2022\u2023\u25E6\-\*]\s", text))

    def _runs_to_inline(self, runs) -> str:
        parts: List[str] = []
        for run in runs:
            text = run.text or ""
            if not text:
                continue
            try:
                bold = run.font.bold
                italic = run.font.italic
            except Exception:
                bold = italic = False
            if bold and italic:
                text = f"***{text}***"
            elif bold:
                text = f"**{text}**"
            elif italic:
                text = f"*{text}*"
            parts.append(text)
        return "".join(parts)

    def _table_to_markdown(self, table) -> str:
        rows: List[List[str]] = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        col_count = max(len(r) for r in rows)
        # Pad rows
        rows = [r + [""] * (col_count - len(r)) for r in rows]

        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join(["---"] * col_count) + " |"
        body_rows = ["| " + " | ".join(r) + " |" for r in rows[1:]]

        return "\n".join([header, separator] + body_rows)

    def _is_image(self, shape) -> bool:
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        except Exception:
            return False

    def _image_to_markdown(self, shape, assets_dir: Path, img_counter: List[int]) -> str:
        try:
            img_counter[0] += 1
            assets_dir.mkdir(parents=True, exist_ok=True)
            img_path = assets_dir / f"img_{img_counter[0]}.png"
            with open(img_path, "wb") as f:
                f.write(shape.image.blob)
            rel_path = f"{assets_dir.name}/img_{img_counter[0]}.png"
            return f"![]({assets_dir.name}/img_{img_counter[0]}.png)"
        except Exception:
            return ""
