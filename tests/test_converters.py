"""
Unit tests for the converter pipeline.

TestLayoutAnalyzer   — spatial grid detection from ShapeRegion lists
TestPPTXConverter    — in-memory python-pptx Presentation → Markdown
TestPDFConverter     — mocked pdfplumber page → Markdown
"""

import io
import textwrap
from pathlib import Path
from unittest.mock import MagicMock, patch, PropertyMock

import pytest

from src.converters.layout_analyzer import LayoutAnalyzer, ShapeRegion, SlideLayout


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def make_region(content: str, x: float, y: float, w: float, h: float) -> ShapeRegion:
    return ShapeRegion(content_md=content, x=x, y=y, w=w, h=h)


# ===========================================================================
# TestLayoutAnalyzer
# ===========================================================================

class TestLayoutAnalyzer:
    def setup_method(self):
        self.analyzer = LayoutAnalyzer()

    # --- analyze() ---

    def test_empty_shapes_returns_single_col(self):
        layout = self.analyzer.analyze([])
        assert layout.n_cols == 1
        assert layout.title == ""
        assert layout.bands == []

    def test_single_shape_no_grid(self):
        shapes = [make_region("Hello world", 0.1, 0.3, 0.8, 0.2)]
        layout = self.analyzer.analyze(shapes)
        assert layout.n_cols == 1

    def test_title_extracted_by_y_threshold(self):
        title = make_region("My Title", 0.1, 0.05, 0.8, 0.1)
        body = make_region("Body text", 0.1, 0.3, 0.8, 0.2)
        layout = self.analyzer.analyze([title, body])
        assert layout.title == "My Title"
        assert len(layout.bands) == 1
        assert layout.bands[0][0].content_md == "Body text"

    def test_two_side_by_side_shapes_gives_two_cols(self):
        left = make_region("Left col", 0.0, 0.3, 0.45, 0.4)
        right = make_region("Right col", 0.55, 0.3, 0.45, 0.4)
        layout = self.analyzer.analyze([left, right])
        assert layout.n_cols == 2

    def test_three_columns_detected(self):
        shapes = [
            make_region("A", 0.0, 0.3, 0.28, 0.4),
            make_region("B", 0.36, 0.3, 0.28, 0.4),
            make_region("C", 0.72, 0.3, 0.28, 0.4),
        ]
        layout = self.analyzer.analyze(shapes)
        assert layout.n_cols == 3

    def test_stacked_shapes_form_single_column(self):
        top = make_region("Top block", 0.1, 0.3, 0.8, 0.2)
        bottom = make_region("Bottom block", 0.1, 0.6, 0.8, 0.2)
        layout = self.analyzer.analyze([top, bottom])
        # Two shapes in two separate vertical bands → max 1 col per band → n_cols == 1
        assert layout.n_cols == 1

    def test_no_title_when_all_shapes_below_threshold(self):
        # All shapes below y=0.25 threshold — no title extracted, both go into bands
        shapes = [
            make_region("First", 0.1, 0.3, 0.8, 0.1),
            make_region("Second", 0.1, 0.5, 0.8, 0.1),
        ]
        layout = self.analyzer.analyze(shapes)
        assert layout.title == ""
        assert len(layout.bands) == 2

    def test_bands_sorted_top_to_bottom(self):
        bottom = make_region("B", 0.1, 0.7, 0.8, 0.1)
        top = make_region("T", 0.1, 0.4, 0.8, 0.1)
        layout = self.analyzer.analyze([bottom, top])
        # Both in separate bands; bands sorted by y
        assert layout.bands[0][0].content_md == "T"
        assert layout.bands[1][0].content_md == "B"

    # --- to_markdown() ---

    def test_single_col_markdown_no_grid_markers(self):
        layout = SlideLayout(
            title="My Slide",
            bands=[[make_region("Some text", 0.1, 0.3, 0.8, 0.2)]],
            n_cols=1,
        )
        md = self.analyzer.to_markdown(layout)
        assert "<!-- $grid" not in md
        assert "# My Slide" in md
        assert "Some text" in md

    def test_two_col_markdown_has_grid_markers(self):
        layout = SlideLayout(
            title="Grid Slide",
            bands=[[
                make_region("Left content", 0.0, 0.3, 0.45, 0.4),
                make_region("Right content", 0.55, 0.3, 0.45, 0.4),
            ]],
            n_cols=2,
        )
        md = self.analyzer.to_markdown(layout)
        assert "<!-- $grid(2) -->" in md
        assert "<!-- $grid/ -->" in md
        assert "Left content" in md
        assert "Right content" in md
        assert "-----" in md  # cell separator

    def test_empty_title_omitted_from_markdown(self):
        layout = SlideLayout(
            title="",
            bands=[[make_region("Just text", 0.1, 0.3, 0.8, 0.2)]],
            n_cols=1,
        )
        md = self.analyzer.to_markdown(layout)
        assert md.startswith("Just text") or "Just text" in md
        assert "#" not in md.split("\n")[0]

    def test_span_detection_wide_shape(self):
        # A shape spanning ~full width in a 2-col layout → span=2
        wide = make_region("Wide content", 0.0, 0.3, 0.95, 0.2)
        span = self.analyzer._compute_span(wide, n_cols=2)
        assert span == 2

    def test_span_detection_half_width_shape(self):
        half = make_region("Half", 0.0, 0.3, 0.45, 0.2)
        span = self.analyzer._compute_span(half, n_cols=2)
        assert span == 1

    def test_grid_cell_comment_emitted_for_spanning_shape(self):
        layout = SlideLayout(
            title="",
            bands=[[
                make_region("Full width", 0.0, 0.3, 0.95, 0.2),
                make_region("Right", 0.55, 0.3, 0.45, 0.2),
            ]],
            n_cols=2,
        )
        md = self.analyzer.to_markdown(layout)
        assert "<!-- $grid-cell(2, 1) -->" in md


# ===========================================================================
# TestPPTXConverter
# ===========================================================================

class TestPPTXConverter:
    """
    Builds minimal in-memory presentations using python-pptx and tests that
    PPTXConverter produces correct Markdown.

    Skipped if python-pptx is not installed.
    """

    @pytest.fixture(autouse=True)
    def require_pptx(self):
        pytest.importorskip("pptx")

    def _blank_prs(self):
        from pptx import Presentation
        return Presentation()

    def _add_text_slide(self, prs, title_text: str, body_text: str):
        from pptx.util import Inches, Pt
        slide_layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        tf = slide.placeholders[1].text_frame
        tf.text = body_text
        return slide

    # --- Basic conversion ---

    def test_convert_returns_string(self, tmp_path):
        from pptx import Presentation
        prs = self._blank_prs()
        self._add_text_slide(prs, "Hello", "World")
        pptx_file = tmp_path / "test.pptx"
        prs.save(str(pptx_file))

        from src.converters.pptx_converter import PPTXConverter
        md = PPTXConverter().convert(pptx_file)
        assert isinstance(md, str)
        assert len(md) > 0

    def test_single_slide_no_separator(self, tmp_path):
        from pptx import Presentation
        prs = self._blank_prs()
        self._add_text_slide(prs, "Only Slide", "Content here")
        pptx_file = tmp_path / "single.pptx"
        prs.save(str(pptx_file))

        from src.converters.pptx_converter import PPTXConverter
        md = PPTXConverter().convert(pptx_file)
        assert "---" not in md or md.count("\n---\n") == 0

    def test_two_slides_produces_separator(self, tmp_path):
        from pptx import Presentation
        prs = self._blank_prs()
        self._add_text_slide(prs, "Slide One", "First")
        self._add_text_slide(prs, "Slide Two", "Second")
        pptx_file = tmp_path / "two.pptx"
        prs.save(str(pptx_file))

        from src.converters.pptx_converter import PPTXConverter
        md = PPTXConverter().convert(pptx_file)
        assert "\n---\n" in md

    def test_title_placeholder_becomes_heading(self, tmp_path):
        from pptx import Presentation
        prs = self._blank_prs()
        self._add_text_slide(prs, "My Title", "Body")
        pptx_file = tmp_path / "title.pptx"
        prs.save(str(pptx_file))

        from src.converters.pptx_converter import PPTXConverter
        md = PPTXConverter().convert(pptx_file)
        assert "My Title" in md

    def test_body_text_present_in_output(self, tmp_path):
        from pptx import Presentation
        prs = self._blank_prs()
        self._add_text_slide(prs, "Title", "Important body text")
        pptx_file = tmp_path / "body.pptx"
        prs.save(str(pptx_file))

        from src.converters.pptx_converter import PPTXConverter
        md = PPTXConverter().convert(pptx_file)
        assert "Important body text" in md

    def test_slide_count_matches(self, tmp_path):
        from pptx import Presentation
        prs = self._blank_prs()
        for i in range(3):
            self._add_text_slide(prs, f"Slide {i+1}", f"Content {i+1}")
        pptx_file = tmp_path / "three.pptx"
        prs.save(str(pptx_file))

        from src.converters.pptx_converter import PPTXConverter
        md = PPTXConverter().convert(pptx_file)
        # 3 slides → 2 separators
        assert md.count("\n---\n") == 2

    # --- Table conversion ---

    def test_table_shape_produces_pipe_table(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        prs = self._blank_prs()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
        rows, cols = 2, 2
        table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2)).table
        table.cell(0, 0).text = "Header A"
        table.cell(0, 1).text = "Header B"
        table.cell(1, 0).text = "Val 1"
        table.cell(1, 1).text = "Val 2"
        pptx_file = tmp_path / "table.pptx"
        prs.save(str(pptx_file))

        from src.converters.pptx_converter import PPTXConverter
        md = PPTXConverter().convert(pptx_file)
        assert "| Header A" in md
        assert "| ---" in md

    # --- PPTXConverter internal helpers (unit-level) ---

    def test_runs_to_inline_bold(self):
        from src.converters.pptx_converter import PPTXConverter
        conv = PPTXConverter()
        run = MagicMock()
        run.text = "bold text"
        run.font.bold = True
        run.font.italic = False
        result = conv._runs_to_inline([run])
        assert result == "**bold text**"

    def test_runs_to_inline_italic(self):
        from src.converters.pptx_converter import PPTXConverter
        conv = PPTXConverter()
        run = MagicMock()
        run.text = "italic"
        run.font.bold = False
        run.font.italic = True
        result = conv._runs_to_inline([run])
        assert result == "*italic*"

    def test_runs_to_inline_plain(self):
        from src.converters.pptx_converter import PPTXConverter
        conv = PPTXConverter()
        run = MagicMock()
        run.text = "plain"
        run.font.bold = False
        run.font.italic = False
        result = conv._runs_to_inline([run])
        assert result == "plain"

    def test_table_to_markdown_structure(self):
        from src.converters.pptx_converter import PPTXConverter
        conv = PPTXConverter()
        table = MagicMock()
        cell_00 = MagicMock(); cell_00.text = "H1"
        cell_01 = MagicMock(); cell_01.text = "H2"
        cell_10 = MagicMock(); cell_10.text = "R1"
        cell_11 = MagicMock(); cell_11.text = "R2"
        row0 = MagicMock(); row0.cells = [cell_00, cell_01]
        row1 = MagicMock(); row1.cells = [cell_10, cell_11]
        table.rows = [row0, row1]
        md = conv._table_to_markdown(table)
        assert "| H1 | H2 |" in md
        assert "| --- | --- |" in md
        assert "| R1 | R2 |" in md

    def test_detect_callout_by_name(self):
        from src.converters.pptx_converter import PPTXConverter
        conv = PPTXConverter()
        shape = MagicMock()
        shape.name = "Info Box 1"
        shape.text_frame.text = "Some text"
        result = conv._detect_callout(shape)
        assert result == "info"

    def test_detect_callout_by_text_start(self):
        from src.converters.pptx_converter import PPTXConverter
        conv = PPTXConverter()
        shape = MagicMock()
        shape.name = "Rectangle"
        shape.text_frame.text = "warning: check this"
        result = conv._detect_callout(shape)
        assert result == "warning"

    def test_detect_callout_none(self):
        from src.converters.pptx_converter import PPTXConverter
        conv = PPTXConverter()
        shape = MagicMock()
        shape.name = "TextBox"
        shape.text_frame.text = "Normal paragraph"
        result = conv._detect_callout(shape)
        assert result is None


# ===========================================================================
# TestPDFConverter
# ===========================================================================

class TestPDFConverter:
    """
    Tests PDFConverter using mocked pdfplumber pages.
    If pdfplumber is available, a smoke test using an in-memory PDF may run.
    """

    @pytest.fixture(autouse=True)
    def require_pdfplumber(self):
        pytest.importorskip("pdfplumber")

    def _make_word(self, text, x0, top, x1=None, bottom=None, size=12, fontname="Helvetica"):
        if x1 is None:
            x1 = x0 + len(text) * 6
        if bottom is None:
            bottom = top + size
        return {
            "text": text, "x0": x0, "top": top,
            "x1": x1, "bottom": bottom,
            "size": size, "fontname": fontname,
        }

    # --- _group_words_into_lines ---

    def test_words_same_top_grouped_into_one_line(self):
        from src.converters.pdf_converter import PDFConverter
        conv = PDFConverter()
        words = [
            self._make_word("Hello", 10, 100),
            self._make_word("World", 60, 100),
        ]
        lines = conv._group_words_into_lines(words)
        assert len(lines) == 1
        assert len(lines[0]) == 2

    def test_words_different_top_grouped_into_two_lines(self):
        from src.converters.pdf_converter import PDFConverter
        conv = PDFConverter()
        words = [
            self._make_word("Line1", 10, 100),
            self._make_word("Line2", 10, 120),
        ]
        lines = conv._group_words_into_lines(words)
        assert len(lines) == 2

    def test_empty_words_returns_empty_lines(self):
        from src.converters.pdf_converter import PDFConverter
        conv = PDFConverter()
        assert conv._group_words_into_lines([]) == []

    # --- _group_lines_into_blocks ---

    def test_lines_close_together_form_one_block(self):
        from src.converters.pdf_converter import PDFConverter
        conv = PDFConverter()
        line1 = [self._make_word("A", 10, 100, bottom=112)]
        line2 = [self._make_word("B", 10, 114, bottom=126)]
        blocks = conv._group_lines_into_blocks([line1, line2])
        assert len(blocks) == 1

    def test_lines_far_apart_form_two_blocks(self):
        from src.converters.pdf_converter import PDFConverter
        conv = PDFConverter()
        line1 = [self._make_word("A", 10, 100, bottom=112)]
        line2 = [self._make_word("B", 10, 150, bottom=162)]
        blocks = conv._group_lines_into_blocks([line1, line2])
        assert len(blocks) == 2

    # --- _classify_line ---

    def test_large_font_classified_as_heading(self):
        from src.converters.pdf_converter import PDFConverter
        conv = PDFConverter()
        block_words = [self._make_word("Title", 10, 10, size=24)]
        result = conv._classify_line("Title", size=24, median_size=12, block_words=block_words)
        assert result.startswith("#")

    def test_normal_font_classified_as_paragraph(self):
        from src.converters.pdf_converter import PDFConverter
        conv = PDFConverter()
        block_words = [self._make_word("Body", 10, 10, size=12)]
        result = conv._classify_line("Body", size=12, median_size=12, block_words=block_words)
        assert result == "Body"

    def test_bullet_character_classified_as_list_item(self):
        from src.converters.pdf_converter import PDFConverter
        conv = PDFConverter()
        block_words = []
        result = conv._classify_line("\u2022 Item text", size=12, median_size=12, block_words=block_words)
        assert result == "- Item text"

    def test_numbered_list_classified_as_list_item(self):
        from src.converters.pdf_converter import PDFConverter
        conv = PDFConverter()
        result = conv._classify_line("1. First item", size=12, median_size=12, block_words=[])
        assert result == "- First item"

    # --- Full page conversion with mocked pdfplumber ---

    def test_convert_mocked_page(self, tmp_path):
        from src.converters.pdf_converter import PDFConverter

        page = MagicMock()
        page.width = 612
        page.height = 792
        page.images = []
        page.extract_words.return_value = [
            self._make_word("Slide", 100, 50, size=28),
            self._make_word("Title", 160, 50, size=28),
            self._make_word("Body", 100, 200, size=12),
            self._make_word("text", 145, 200, size=12),
        ]

        mock_pdf = MagicMock()
        mock_pdf.__enter__ = MagicMock(return_value=mock_pdf)
        mock_pdf.__exit__ = MagicMock(return_value=False)
        mock_pdf.pages = [page]

        dummy_pdf = tmp_path / "dummy.pdf"
        dummy_pdf.write_bytes(b"%PDF-1.4")

        with patch("pdfplumber.open", return_value=mock_pdf):
            md = PDFConverter().convert(dummy_pdf)

        assert isinstance(md, str)
        assert "Slide" in md or "Title" in md

    def test_convert_two_pages_produces_separator(self, tmp_path):
        from src.converters.pdf_converter import PDFConverter

        def make_page(text):
            page = MagicMock()
            page.width = 612
            page.height = 792
            page.images = []
            page.extract_words.return_value = [
                self._make_word(text, 100, 100, size=12),
            ]
            return page

        mock_pdf = MagicMock()
        mock_pdf.__enter__ = MagicMock(return_value=mock_pdf)
        mock_pdf.__exit__ = MagicMock(return_value=False)
        mock_pdf.pages = [make_page("PageOne"), make_page("PageTwo")]

        dummy_pdf = tmp_path / "dummy2.pdf"
        dummy_pdf.write_bytes(b"%PDF-1.4")

        with patch("pdfplumber.open", return_value=mock_pdf):
            md = PDFConverter().convert(dummy_pdf)

        assert "\n---\n" in md or "---" in md


# ===========================================================================
# TestConverterCLI (integration smoke)
# ===========================================================================

class TestConverterCLI:
    def test_import_convert_module(self):
        import convert  # noqa: F401

    def test_pptx_converter_import(self):
        from src.converters.pptx_converter import PPTXConverter
        assert callable(PPTXConverter)

    def test_pdf_converter_import(self):
        from src.converters.pdf_converter import PDFConverter
        assert callable(PDFConverter)

    def test_base_converter_is_abstract(self):
        from src.converters.base import PresentationConverter
        import inspect
        assert inspect.isabstract(PresentationConverter)

    def test_layout_analyzer_import(self):
        from src.converters.layout_analyzer import LayoutAnalyzer, ShapeRegion, SlideLayout
        assert callable(LayoutAnalyzer)
