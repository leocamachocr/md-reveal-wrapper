"""
Unit tests for presentation_importer and grid_builder.

TestGridBuilder           — GridBuilder.image_block / n_col / cell_span
TestPresentationImporter  — factory (PresentationImporter.for_file)
TestImportResult          — ImportResult dataclass fields
TestPPTXImporter          — in-memory python-pptx → Markdown (skipped without lib)
TestPDFImporter           — mocked pdfplumber page → Markdown (skipped without lib)
"""

import io
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

from src.grid_builder import GridBuilder
from src.presentation_importer import (
    ImportResult,
    PresentationImporter,
    PPTXImporter,
    PDFImporter,
)


# ===========================================================================
# TestGridBuilder
# ===========================================================================

class TestGridBuilder:
    def setup_method(self):
        self.gb = GridBuilder()

    # --- image_block ---

    def test_wide_image_returns_plain_img_tag(self):
        md = self.gb.image_block("images/img.png", x_frac=0.05, w_frac=0.90)
        assert md == "![](images/img.png)"
        assert "grid" not in md

    def test_centred_image_returns_plain_img_tag(self):
        md = self.gb.image_block("images/img.png", x_frac=0.50, w_frac=0.30)
        assert md == "![](images/img.png)"
        assert "grid" not in md

    def test_left_image_wrapped_in_left_column(self):
        md = self.gb.image_block("images/img.png", x_frac=0.05, w_frac=0.35)
        assert "<!-- $grid(2) -->" in md
        assert "<!-- $grid/ -->" in md
        lines = md.split("\n")
        # The img tag appears before the cell separator
        img_idx = next(i for i, l in enumerate(lines) if "![" in l)
        sep_idx = next(i for i, l in enumerate(lines) if l.strip() == "-----")
        assert img_idx < sep_idx

    def test_right_image_wrapped_in_right_column(self):
        md = self.gb.image_block("images/img.png", x_frac=0.70, w_frac=0.25)
        assert "<!-- $grid(2) -->" in md
        lines = md.split("\n")
        img_idx = next(i for i, l in enumerate(lines) if "![" in l)
        sep_idx = next(i for i, l in enumerate(lines) if l.strip() == "-----")
        # img appears AFTER the separator (right column)
        assert img_idx > sep_idx

    def test_boundary_x_040_is_centred(self):
        md = self.gb.image_block("images/img.png", x_frac=0.40, w_frac=0.20)
        assert "grid" not in md

    def test_boundary_x_060_is_centred(self):
        md = self.gb.image_block("images/img.png", x_frac=0.60, w_frac=0.20)
        assert "grid" not in md

    def test_boundary_w_060_is_full_width(self):
        md = self.gb.image_block("images/img.png", x_frac=0.05, w_frac=0.60)
        assert "grid" not in md

    # --- n_col ---

    def test_n_col_two_cells_structure(self):
        md = self.gb.n_col(["Left", "Right"], n=2)
        assert "<!-- $grid(2) -->" in md
        assert "<!-- $grid/ -->" in md
        assert "-----" in md
        assert "Left" in md
        assert "Right" in md

    def test_n_col_three_cells_has_two_separators(self):
        md = self.gb.n_col(["A", "B", "C"], n=3)
        assert "<!-- $grid(3) -->" in md
        assert md.count("-----") == 2

    def test_n_col_empty_cell_preserved(self):
        md = self.gb.n_col(["Content", ""], n=2)
        assert "Content" in md
        assert "-----" in md

    def test_n_col_cell_order_is_preserved(self):
        md = self.gb.n_col(["First", "Second"], n=2)
        first_pos = md.index("First")
        sep_pos = md.index("-----")
        second_pos = md.index("Second")
        assert first_pos < sep_pos < second_pos

    # --- cell_span ---

    def test_cell_span_emits_directive(self):
        md = self.gb.cell_span("Wide content", col_span=2)
        assert "<!-- $grid-cell(2, 1) -->" in md
        assert "Wide content" in md

    def test_cell_span_custom_row(self):
        md = self.gb.cell_span("Tall cell", col_span=1, row_span=3)
        assert "<!-- $grid-cell(1, 3) -->" in md

    def test_cell_span_strips_extra_whitespace(self):
        md = self.gb.cell_span("  trimmed  ", col_span=2)
        assert "trimmed" in md
        assert "  trimmed  " not in md


# ===========================================================================
# TestPresentationImporter (factory)
# ===========================================================================

class TestPresentationImporter:
    def test_pptx_extension_returns_pptx_importer(self, tmp_path):
        f = tmp_path / "test.pptx"
        f.touch()
        importer = PresentationImporter.for_file(f)
        assert isinstance(importer, PPTXImporter)

    def test_ppt_extension_returns_pptx_importer(self, tmp_path):
        f = tmp_path / "old.ppt"
        f.touch()
        importer = PresentationImporter.for_file(f)
        assert isinstance(importer, PPTXImporter)

    def test_pdf_extension_returns_pdf_importer(self, tmp_path):
        f = tmp_path / "doc.pdf"
        f.touch()
        importer = PresentationImporter.for_file(f)
        assert isinstance(importer, PDFImporter)

    def test_unsupported_extension_raises_value_error(self, tmp_path):
        f = tmp_path / "slides.key"
        f.touch()
        with pytest.raises(ValueError, match="Unsupported format"):
            PresentationImporter.for_file(f)

    def test_case_insensitive_extension(self, tmp_path):
        f = tmp_path / "UPPER.PPTX"
        f.touch()
        importer = PresentationImporter.for_file(f)
        assert isinstance(importer, PPTXImporter)


# ===========================================================================
# TestImportResult
# ===========================================================================

class TestImportResult:
    def test_fields_are_accessible(self, tmp_path):
        result = ImportResult(
            md_content="# Hello",
            slide_count=3,
            image_count=2,
            output_path=tmp_path / "out.md",
        )
        assert result.slide_count == 3
        assert result.image_count == 2
        assert result.md_content == "# Hello"

    def test_output_path_is_path_type(self, tmp_path):
        result = ImportResult(
            md_content="",
            slide_count=0,
            image_count=0,
            output_path=tmp_path / "out.md",
        )
        assert isinstance(result.output_path, Path)


# ===========================================================================
# TestPPTXImporter
# ===========================================================================

class TestPPTXImporter:
    """
    Builds minimal in-memory presentations with python-pptx and tests that
    PPTXImporter writes correct Markdown and returns correct stats.

    Skipped if python-pptx is not installed.
    """

    @pytest.fixture(autouse=True)
    def require_pptx(self):
        pytest.importorskip("pptx")

    def _make_pptx(self, tmp_path, slides_data):
        """
        slides_data: list of (title, body) tuples.
        Returns path to the saved .pptx file.
        """
        from pptx import Presentation
        prs = Presentation()
        layout = prs.slide_layouts[1]  # Title + Content
        for title_text, body_text in slides_data:
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title_text
            slide.placeholders[1].text_frame.text = body_text
        path = tmp_path / "test.pptx"
        prs.save(str(path))
        return path

    def test_import_returns_import_result(self, tmp_path):
        path = self._make_pptx(tmp_path, [("Title", "Body")])
        result = PPTXImporter().import_file(path)
        assert isinstance(result, ImportResult)

    def test_slide_count_matches(self, tmp_path):
        path = self._make_pptx(tmp_path, [
            ("S1", "Body1"), ("S2", "Body2"), ("S3", "Body3")
        ])
        result = PPTXImporter().import_file(path)
        assert result.slide_count == 3

    def test_output_file_written(self, tmp_path):
        path = self._make_pptx(tmp_path, [("Hello", "World")])
        result = PPTXImporter().import_file(path)
        assert result.output_path.is_file()
        assert result.output_path.suffix == ".md"

    def test_output_file_content_matches_result_md(self, tmp_path):
        path = self._make_pptx(tmp_path, [("Hello", "World")])
        result = PPTXImporter().import_file(path)
        on_disk = result.output_path.read_text(encoding="utf-8")
        assert on_disk == result.md_content

    def test_title_text_present_in_output(self, tmp_path):
        path = self._make_pptx(tmp_path, [("My Heading", "Details")])
        result = PPTXImporter().import_file(path)
        assert "My Heading" in result.md_content

    def test_body_text_present_in_output(self, tmp_path):
        path = self._make_pptx(tmp_path, [("T", "Important body content")])
        result = PPTXImporter().import_file(path)
        assert "Important body content" in result.md_content

    def test_multiple_slides_have_separator(self, tmp_path):
        path = self._make_pptx(tmp_path, [("S1", "A"), ("S2", "B")])
        result = PPTXImporter().import_file(path)
        assert "---" in result.md_content

    def test_custom_output_path(self, tmp_path):
        path = self._make_pptx(tmp_path, [("X", "Y")])
        out = tmp_path / "custom_output.md"
        result = PPTXImporter().import_file(path, output_path=out)
        assert result.output_path == out
        assert out.is_file()

    def test_notes_appended_when_present(self, tmp_path):
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide with notes"
        slide.placeholders[1].text_frame.text = "Content"
        # Add notes
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = "Speaker note text"
        path = tmp_path / "noted.pptx"
        prs.save(str(path))

        result = PPTXImporter().import_file(path)
        assert "<!-- notes: Speaker note text -->" in result.md_content

    def test_zero_images_when_no_pictures(self, tmp_path):
        path = self._make_pptx(tmp_path, [("T", "No images here")])
        result = PPTXImporter().import_file(path)
        assert result.image_count == 0

    def test_image_saved_to_images_dir(self, tmp_path):
        """Smoke test: import a file with a picture shape saves to images/."""
        from pptx import Presentation
        from pptx.util import Inches
        import struct, zlib

        # Build a minimal 1×1 white PNG in memory.
        def _minimal_png():
            def chunk(tag, data):
                c = zlib.crc32(tag + data) & 0xFFFFFFFF
                return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", c)
            sig = b"\x89PNG\r\n\x1a\n"
            ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
            idat = chunk(b"IDAT", zlib.compress(b"\x00\xff\xff\xff"))
            iend = chunk(b"IEND", b"")
            return sig + ihdr + idat + iend

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
        png_bytes = _minimal_png()
        img_stream = io.BytesIO(png_bytes)
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1), Inches(2), Inches(2))
        path = tmp_path / "with_image.pptx"
        prs.save(str(path))

        result = PPTXImporter().import_file(path)
        assert result.image_count == 1
        images_dir = tmp_path / "images"
        assert images_dir.is_dir()
        saved_files = list(images_dir.glob("slide_1_img_1.png"))
        assert len(saved_files) == 1

    def test_image_reference_uses_slide_aware_name(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        import struct, zlib

        def _minimal_png():
            def chunk(tag, data):
                c = zlib.crc32(tag + data) & 0xFFFFFFFF
                return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", c)
            sig = b"\x89PNG\r\n\x1a\n"
            ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
            idat = chunk(b"IDAT", zlib.compress(b"\x00\xff\xff\xff"))
            iend = chunk(b"IEND", b"")
            return sig + ihdr + idat + iend

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        img_stream = io.BytesIO(_minimal_png())
        slide.shapes.add_picture(img_stream, Inches(1), Inches(1), Inches(2), Inches(2))
        path = tmp_path / "img_name.pptx"
        prs.save(str(path))

        result = PPTXImporter().import_file(path)
        assert "images/slide_1_img_1.png" in result.md_content


# ===========================================================================
# TestPDFImporter
# ===========================================================================

class TestPDFImporter:
    """
    Tests PDFImporter using mocked pdfplumber pages.
    Skipped if pdfplumber is not installed.
    """

    @pytest.fixture(autouse=True)
    def require_pdfplumber(self):
        pytest.importorskip("pdfplumber")

    def _make_word(self, text, x0, top, size=12):
        x1 = x0 + len(text) * 6
        bottom = top + size
        return {
            "text": text, "x0": x0, "top": top,
            "x1": x1, "bottom": bottom,
            "size": size, "fontname": "Helvetica",
        }

    def _mock_page(self, words, images=None):
        page = MagicMock()
        page.width = 612
        page.height = 792
        page.images = images or []
        page.extract_words.return_value = words
        return page

    def _mock_pdf(self, pages):
        pdf = MagicMock()
        pdf.__enter__ = MagicMock(return_value=pdf)
        pdf.__exit__ = MagicMock(return_value=False)
        pdf.pages = pages
        return pdf

    def test_import_returns_import_result(self, tmp_path):
        page = self._mock_page([self._make_word("Hello", 100, 100)])
        dummy = tmp_path / "test.pdf"
        dummy.write_bytes(b"%PDF-1.4")

        with patch("pdfplumber.open", return_value=self._mock_pdf([page])):
            result = PDFImporter().import_file(dummy)

        assert isinstance(result, ImportResult)

    def test_slide_count_matches_page_count(self, tmp_path):
        pages = [
            self._mock_page([self._make_word("Page1", 100, 100)]),
            self._mock_page([self._make_word("Page2", 100, 100)]),
        ]
        dummy = tmp_path / "two.pdf"
        dummy.write_bytes(b"%PDF-1.4")

        with patch("pdfplumber.open", return_value=self._mock_pdf(pages)):
            result = PDFImporter().import_file(dummy)

        assert result.slide_count == 2

    def test_text_present_in_output(self, tmp_path):
        page = self._mock_page([self._make_word("Visible", 100, 100)])
        dummy = tmp_path / "text.pdf"
        dummy.write_bytes(b"%PDF-1.4")

        with patch("pdfplumber.open", return_value=self._mock_pdf([page])):
            result = PDFImporter().import_file(dummy)

        assert "Visible" in result.md_content

    def test_multiple_pages_have_separator(self, tmp_path):
        pages = [
            self._mock_page([self._make_word("P1", 100, 100)]),
            self._mock_page([self._make_word("P2", 100, 100)]),
        ]
        dummy = tmp_path / "sep.pdf"
        dummy.write_bytes(b"%PDF-1.4")

        with patch("pdfplumber.open", return_value=self._mock_pdf(pages)):
            result = PDFImporter().import_file(dummy)

        assert "---" in result.md_content

    def test_output_file_written(self, tmp_path):
        page = self._mock_page([self._make_word("Text", 100, 100)])
        dummy = tmp_path / "out.pdf"
        dummy.write_bytes(b"%PDF-1.4")

        with patch("pdfplumber.open", return_value=self._mock_pdf([page])):
            result = PDFImporter().import_file(dummy)

        assert result.output_path.is_file()
        assert result.output_path.suffix == ".md"

    def test_custom_output_path(self, tmp_path):
        page = self._mock_page([self._make_word("X", 100, 100)])
        dummy = tmp_path / "src.pdf"
        dummy.write_bytes(b"%PDF-1.4")
        out = tmp_path / "my_output.md"

        with patch("pdfplumber.open", return_value=self._mock_pdf([page])):
            result = PDFImporter().import_file(dummy, output_path=out)

        assert result.output_path == out
        assert out.is_file()

    def test_zero_images_when_no_images_on_page(self, tmp_path):
        page = self._mock_page([self._make_word("NoImg", 100, 100)], images=[])
        dummy = tmp_path / "noimg.pdf"
        dummy.write_bytes(b"%PDF-1.4")

        with patch("pdfplumber.open", return_value=self._mock_pdf([page])):
            result = PDFImporter().import_file(dummy)

        assert result.image_count == 0

    def test_output_content_matches_on_disk(self, tmp_path):
        page = self._mock_page([self._make_word("Consistent", 100, 100)])
        dummy = tmp_path / "cons.pdf"
        dummy.write_bytes(b"%PDF-1.4")

        with patch("pdfplumber.open", return_value=self._mock_pdf([page])):
            result = PDFImporter().import_file(dummy)

        on_disk = result.output_path.read_text(encoding="utf-8")
        assert on_disk == result.md_content
